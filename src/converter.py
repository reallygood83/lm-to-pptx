"""
NotebookLM PDF to PPTX Converter
Main converter class with AI-powered speaker notes generation
"""

import os
import shutil
from pathlib import Path
from typing import Optional, Union, List
from PIL import Image

try:
    from pdf2image import convert_from_path
except ImportError:
    convert_from_path = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    Presentation = None

from .ai_providers import (
    AIProvider,
    GeminiProvider,
    OpenAIProvider,
    AnthropicProvider,
    GrokProvider
)


class NotebookLMToPPTX:
    """
    NotebookLM PDF를 편집 가능한 PPTX로 변환하는 메인 컨버터.

    주요 기능:
    - PDF 슬라이드를 고품질 이미지로 변환
    - 16:9 PowerPoint 표준 규격에 맞게 풀슬라이드 이미지 삽입
    - AI Vision API를 활용한 스피커 노트 자동 생성
    - 맥락 자료(Context Materials)를 활용한 노트 품질 향상
    """

    # PowerPoint 16:9 표준 크기 (인치)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # 지원하는 AI 프로바이더
    PROVIDERS = {
        'gemini': GeminiProvider,
        'openai': OpenAIProvider,
        'anthropic': AnthropicProvider,
        'claude': AnthropicProvider,  # alias
        'grok': GrokProvider,
        'xai': GrokProvider,  # alias
    }

    def __init__(
        self,
        provider: str = 'gemini',
        api_key: Optional[str] = None,
        model: Optional[str] = None,
        dpi: int = 144
    ):
        """
        컨버터 초기화.

        Args:
            provider: AI 프로바이더 이름 ('gemini', 'openai', 'anthropic', 'grok')
            api_key: API 키 (환경변수에서 자동 로드 가능)
            model: 사용할 모델명 (None이면 기본값 사용)
            dpi: PDF 렌더링 해상도 (기본: 144 DPI)
        """
        self._check_dependencies()

        self.dpi = dpi
        self.provider_name = provider.lower()

        # AI 프로바이더 설정
        if self.provider_name not in self.PROVIDERS:
            raise ValueError(
                f"지원하지 않는 프로바이더: {provider}. "
                f"사용 가능: {list(self.PROVIDERS.keys())}"
            )

        # API 키 설정 (환경변수 또는 직접 전달)
        self.api_key = api_key or self._get_api_key_from_env()

        # 프로바이더 인스턴스 생성
        provider_class = self.PROVIDERS[self.provider_name]
        if model:
            self.ai_provider = provider_class(self.api_key, model)
        else:
            self.ai_provider = provider_class(self.api_key)

    def _check_dependencies(self):
        """필수 의존성 및 외부 도구 확인."""
        if convert_from_path is None:
            raise ImportError(
                "pdf2image 패키지가 필요합니다. "
                "설치: pip install pdf2image"
            )

        if Presentation is None:
            raise ImportError(
                "python-pptx 패키지가 필요합니다. "
                "설치: pip install python-pptx"
            )

        # Poppler (pdftoppm) 의존성 확인
        if not shutil.which("pdftoppm") and not shutil.which("pdftocairo"):
            import platform
            os_name = platform.system()
            
            error_msg = (
                "\n" + "="*50 + "\n"
                "🚨 Poppler를 찾을 수 없습니다! (PDF 변환 필수 도구)\n"
                "--------------------------------------------------\n"
                f"현재 운영체제: {os_name}\n\n"
            )
            
            if os_name == "Darwin":  # macOS
                error_msg += "👉 설치 방법: brew install poppler\n"
            elif os_name == "Windows":
                error_msg += (
                    "👉 설치 방법:\n"
                    "1. https://github.com/oschwartz10612/poppler-windows/releases/ 에서 최신 bin.7z 다운로드\n"
                    "2. 압축 해제 후 'bin' 폴더 경로를 시스템 환경변수 'Path'에 추가\n"
                    "3. 또는 'nb2pptx --ui' 명령어로 앱을 켠 뒤 사이드바의 [다운로드 페이지 열기] 버튼 클릭!\n"
                )
            else:
                error_msg += "👉 설치 방법: sudo apt-get install poppler-utils (Ubuntu/Debian)\n"
            
            error_msg += "="*50 + "\n"
            raise RuntimeError(error_msg)

    def _get_api_key_from_env(self) -> str:
        """환경변수에서 API 키 가져오기."""
        env_vars = {
            'gemini': 'GOOGLE_API_KEY',
            'openai': 'OPENAI_API_KEY',
            'anthropic': 'ANTHROPIC_API_KEY',
            'claude': 'ANTHROPIC_API_KEY',
            'grok': 'XAI_API_KEY',
            'xai': 'XAI_API_KEY',
        }

        env_var = env_vars.get(self.provider_name)
        api_key = os.environ.get(env_var, '')

        if not api_key:
            raise ValueError(
                f"API 키가 필요합니다. "
                f"환경변수 {env_var}를 설정하거나 api_key 파라미터를 전달하세요."
            )

        return api_key

    def load_context_materials(
        self,
        paths: Union[str, Path, List[Union[str, Path]]]
    ) -> str:
        """
        맥락 자료 로드.

        Args:
            paths: 파일 경로 또는 경로 리스트 (.txt, .md, .pdf 지원)

        Returns:
            결합된 맥락 텍스트
        """
        if isinstance(paths, (str, Path)):
            paths = [paths]

        context_parts = []

        for path in paths:
            path = Path(path)

            if not path.exists():
                print(f"⚠️ 파일을 찾을 수 없음: {path}")
                continue

            suffix = path.suffix.lower()

            if suffix in ['.txt', '.md']:
                with open(path, 'r', encoding='utf-8') as f:
                    content = f.read()
                context_parts.append(f"--- {path.name} ---\n{content}")

            elif suffix == '.pdf':
                # PDF에서 텍스트 추출 (간단한 방법)
                try:
                    import fitz  # PyMuPDF
                    doc = fitz.open(path)
                    text_parts = []
                    for page in doc:
                        text_parts.append(page.get_text())
                    doc.close()
                    content = '\n'.join(text_parts)
                    context_parts.append(f"--- {path.name} ---\n{content}")
                except ImportError:
                    print(f"⚠️ PDF 텍스트 추출을 위해 PyMuPDF가 필요합니다: pip install pymupdf")
                except Exception as e:
                    print(f"⚠️ PDF 읽기 실패 ({path}): {e}")
            else:
                print(f"⚠️ 지원하지 않는 파일 형식: {suffix}")

        return '\n\n'.join(context_parts)

    def convert_pdf_to_images(self, pdf_path: Union[str, Path]) -> List[Image.Image]:
        """
        PDF를 이미지 리스트로 변환.

        Args:
            pdf_path: PDF 파일 경로

        Returns:
            PIL Image 리스트
        """
        pdf_path = Path(pdf_path)

        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF 파일을 찾을 수 없습니다: {pdf_path}")

        print(f"📄 PDF 로딩 중: {pdf_path.name}")

        images = convert_from_path(
            str(pdf_path),
            dpi=self.dpi,
            fmt='png'
        )

        print(f"✅ {len(images)}개 슬라이드 변환 완료")

        return images

    def create_pptx(
        self,
        images: List[Image.Image],
        output_path: Union[str, Path],
        context: Optional[str] = None,
        generate_notes: bool = True,
        progress_callback: Optional[callable] = None
    ) -> Path:
        """
        이미지 리스트로 PPTX 생성.

        Args:
            images: PIL Image 리스트
            output_path: 출력 PPTX 파일 경로
            context: 스피커 노트 생성용 맥락 자료
            generate_notes: AI 스피커 노트 생성 여부
            progress_callback: 진행 상황 콜백 함수 (current, total)

        Returns:
            생성된 PPTX 파일 경로
        """
        output_path = Path(output_path)

        # 새 프레젠테이션 생성
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT

        # 빈 레이아웃 가져오기
        blank_layout = prs.slide_layouts[6]  # 빈 슬라이드

        total = len(images)

        for idx, image in enumerate(images, 1):
            slide_num = idx

            if progress_callback:
                progress_callback(idx, total)
            else:
                print(f"🔄 슬라이드 {slide_num}/{total} 처리 중...")

            # 슬라이드 추가
            slide = prs.slides.add_slide(blank_layout)

            # 이미지를 임시 파일로 저장
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                image.save(tmp.name, 'PNG')
                tmp_path = tmp.name

            try:
                # 풀슬라이드 이미지 삽입
                slide.shapes.add_picture(
                    tmp_path,
                    Inches(0),
                    Inches(0),
                    width=self.SLIDE_WIDTH,
                    height=self.SLIDE_HEIGHT
                )
            finally:
                # 임시 파일 삭제
                os.unlink(tmp_path)

            # AI 스피커 노트 생성
            if generate_notes:
                try:
                    print(f"  🤖 AI 스피커 노트 생성 중... ({self.provider_name})")
                    notes = self.ai_provider.analyze_slide(image, context)

                    # 노트 슬라이드에 추가
                    notes_slide = slide.notes_slide
                    notes_frame = notes_slide.notes_text_frame
                    notes_frame.text = notes

                    print(f"  ✅ 스피커 노트 생성 완료")

                except Exception as e:
                    print(f"  ⚠️ 스피커 노트 생성 실패: {e}")

        # PPTX 저장
        prs.save(str(output_path))
        print(f"\n🎉 PPTX 저장 완료: {output_path}")

        return output_path

    def convert(
        self,
        pdf_path: Union[str, Path],
        output_path: Optional[Union[str, Path]] = None,
        context_paths: Optional[Union[str, Path, List[Union[str, Path]]]] = None,
        generate_notes: bool = True,
        progress_callback: Optional[callable] = None
    ) -> Path:
        """
        PDF를 PPTX로 변환 (메인 메서드).

        Args:
            pdf_path: 입력 PDF 파일 경로
            output_path: 출력 PPTX 파일 경로 (None이면 자동 생성)
            context_paths: 맥락 자료 파일 경로
            generate_notes: AI 스피커 노트 생성 여부
            progress_callback: 진행 상황 콜백 함수

        Returns:
            생성된 PPTX 파일 경로
        """
        pdf_path = Path(pdf_path)

        # 출력 경로 자동 설정
        if output_path is None:
            output_path = pdf_path.with_suffix('.pptx')
        else:
            output_path = Path(output_path)

        print(f"\n{'='*50}")
        print(f"📊 NotebookLM PDF → PPTX 변환기 v0.4.0")
        print(f"   by 배움의 달인")
        print(f"{'='*50}")
        print(f"\n📥 입력: {pdf_path}")
        print(f"📤 출력: {output_path}")
        print(f"🤖 AI: {self.provider_name} ({self.ai_provider.model})")

        # 맥락 자료 로드
        context = None
        if context_paths:
            print(f"\n📚 맥락 자료 로딩 중...")
            context = self.load_context_materials(context_paths)
            if context:
                print(f"✅ 맥락 자료 로드 완료 ({len(context)} 문자)")

        # PDF → 이미지 변환
        print(f"\n🔄 PDF 변환 중...")
        images = self.convert_pdf_to_images(pdf_path)

        # PPTX 생성
        print(f"\n🎨 PPTX 생성 중...")
        result_path = self.create_pptx(
            images,
            output_path,
            context=context,
            generate_notes=generate_notes,
            progress_callback=progress_callback
        )

        print(f"\n{'='*50}")
        print(f"✨ 변환 완료!")
        print(f"{'='*50}\n")

        return result_path


def quick_convert(
    pdf_path: str,
    provider: str = 'gemini',
    api_key: Optional[str] = None,
    context_paths: Optional[List[str]] = None
) -> Path:
    """
    빠른 변환을 위한 헬퍼 함수.

    Args:
        pdf_path: PDF 파일 경로
        provider: AI 프로바이더 ('gemini', 'openai', 'anthropic', 'grok')
        api_key: API 키 (환경변수 사용 시 생략 가능)
        context_paths: 맥락 자료 경로 리스트

    Returns:
        생성된 PPTX 파일 경로

    Example:
        >>> from src.converter import quick_convert
        >>> quick_convert("slides.pdf", provider="gemini")
    """
    converter = NotebookLMToPPTX(
        provider=provider,
        api_key=api_key
    )

    return converter.convert(
        pdf_path,
        context_paths=context_paths
    )
